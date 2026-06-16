#!/usr/bin/env python3
"""Build an EE deck by cloning the genuine master template's layouts.

The point — colour, layout, the left-side "[" bracket, the genuine EE logo, and
Lexend are all **INHERITED** from the master's layouts/master, never improvised.
Each new slide is added from a real master layout and its placeholders are filled;
nothing is drawn from scratch. This is what makes the output genuinely on-brand
rather than a convincing fake (the Phase-4 failure mode).

Inputs:
  --master    ee-master-2026.pptx  (the bundled asset; resolved by default from
              $CLAUDE_PLUGIN_ROOT/assets/slides or the repo dev path)
  --manifest  assets.manifest.json (layout index map; resolved next to --master)
  --spec      slides.json: {"slides":[{"layout":"title","fields":{"0":"...","2":"..."}}, ...]}
              layout ∈ manifest.layouts (title|section|content|blank|...);
              fields keys = placeholder idx, values = text.
  --out       output .pptx

Guardrails:
  - If the master or manifest can't be resolved, STOP with an error — never build
    an approximation. The genuine logo is inherited from the layouts; this script
    NEVER draws or re-letters a logo.
  - After saving, it verifies the embedded fonts survived (python-pptx preserves
    them; the check is a defensive guard) and that the genuine logo media is present.

Usage:
  python3 build_from_master.py --spec slides.json --out deck.pptx
  python3 build_from_master.py --master M.pptx --manifest A.json --spec S.json --out O.pptx
"""
import argparse
import json
import os
import re
import sys
import zipfile

HERE = os.path.dirname(os.path.abspath(__file__))
BRAND_FONT = "Lexend"


def die(msg):
    print(f"ERROR: {msg}", file=sys.stderr)
    sys.exit(1)


def default_master():
    """Resolve the bundled master: plugin root first, then the repo dev path."""
    root = os.environ.get("CLAUDE_PLUGIN_ROOT")
    candidates = []
    if root:
        candidates.append(os.path.join(root, "assets", "slides", "ee-master-2026.pptx"))
    # repo dev path: skills/create-presentation/scripts -> repo root -> assets/slides
    candidates.append(os.path.abspath(os.path.join(HERE, "..", "..", "..", "assets", "slides", "ee-master-2026.pptx")))
    for c in candidates:
        if os.path.exists(c):
            return c
    die("could not resolve the bundled master (ee-master-2026.pptx). "
        "Set CLAUDE_PLUGIN_ROOT or pass --master. Do NOT build an approximation — stop and flag the missing asset.")


def fill_placeholder(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            for para in ph.text_frame.paragraphs:   # belt-and-braces: force Lexend on the runs
                for run in para.runs:
                    run.font.name = BRAND_FONT
            return True
    return False


def main():
    try:
        from pptx import Presentation
    except ImportError:
        die("python-pptx is required: pip install python-pptx")

    ap = argparse.ArgumentParser()
    ap.add_argument("--master", default=None)
    ap.add_argument("--manifest", default=None)
    ap.add_argument("--spec", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    master = args.master or default_master()
    if not os.path.exists(master):
        die(f"master not found: {master} — stop and flag the missing asset, do not recreate it.")
    manifest_path = args.manifest or os.path.join(os.path.dirname(master), "assets.manifest.json")
    if not os.path.exists(manifest_path):
        die(f"manifest not found: {manifest_path}")
    if not os.path.exists(args.spec):
        die(f"spec not found: {args.spec}")

    manifest = json.load(open(manifest_path))
    layout_map = manifest["layouts"]
    spec = json.load(open(args.spec))

    prs = Presentation(master)
    layouts = list(prs.slide_masters[0].slide_layouts)

    # start from a clean slide list (the master ships 0 slides, but be defensive)
    for sid in list(prs.slides._sldIdLst):
        prs.slides._sldIdLst.remove(sid)

    built = 0
    for i, slide_spec in enumerate(spec["slides"]):
        key = slide_spec["layout"]
        if key not in layout_map or "index" not in layout_map[key]:
            die(f"slide {i}: unknown layout '{key}' (manifest layouts: {sorted(layout_map)})")
        layout = layouts[layout_map[key]["index"]]
        slide = prs.slides.add_slide(layout)
        for idx_str, text in (slide_spec.get("fields") or {}).items():
            idx = int(idx_str)
            if not fill_placeholder(slide, idx, text):
                print(f"  warn: slide {i} ({key}) has no placeholder idx={idx}", file=sys.stderr)
        built += 1

    prs.save(args.out)

    # --- post-save guard: genuine fonts + logo survived ---
    with zipfile.ZipFile(args.out) as z:
        pres_xml = z.read("ppt/presentation.xml").decode("utf-8")
        names = z.namelist()
    embedded = set(re.findall(r'<p:font typeface="([^"]+)"', pres_xml))
    n_fonts = len([n for n in names if n.startswith("ppt/fonts/")])
    logo_present = any(n.startswith("ppt/media/") for n in names)  # logo media inherited via layouts
    warnings = []
    for want in ("Lexend", "Lora", "JetBrains Mono"):
        if want not in embedded:
            warnings.append(f"embedded font lost in build: {want}")
    if not logo_present:
        warnings.append("no inherited media (logo) in output — check the master")

    print(f"wrote {args.out}: {built} slides, {n_fonts} font parts, embedded={sorted(embedded)}")
    if warnings:
        # Defensive: python-pptx preserves embedded fonts, so this should not fire.
        for w in warnings:
            print(f"  WARN: {w}", file=sys.stderr)
    else:
        print("  guard: genuine fonts + logo inherited OK")


if __name__ == "__main__":
    main()
