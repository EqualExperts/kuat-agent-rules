#!/usr/bin/env python3
"""inspect-layouts — enumerate the master's slide layouts to seed manifest labelling.

The studio master ships 65 layouts; only 4 are labelled in assets.manifest.json. Full
layout-map curation (label every layout by purpose, prune unused) is the `curate-layouts`
job. This helper does the mechanical half: it lists every layout with its index (the
python-pptx slide_layouts index — the value the manifest records), name, and placeholder
structure, so a contributor can label them without guessing indices.

It does NOT render or visually judge a layout (no renderer in this environment) — that
visual pass is the contributor step. It only reports structure.

Resolves the master from $CLAUDE_PLUGIN_ROOT/assets/slides (installed) or the repo
assets/slides (dev). Needs python-pptx (`pip install python-pptx`).

Usage:
  python3 skills/scripts/inspect-layouts.py            # human table
  python3 skills/scripts/inspect-layouts.py --json     # JSON seed for the manifest
"""
import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.exit("python-pptx not installed — run: pip install python-pptx")

HERE = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.dirname(os.path.dirname(HERE))


def resolve_master():
    candidates = []
    root = os.environ.get("CLAUDE_PLUGIN_ROOT")
    if root:
        candidates.append(os.path.join(root, "assets", "slides", "ee-master-2026.pptx"))
    candidates.append(os.path.join(REPO_ROOT, "assets", "slides", "ee-master-2026.pptx"))
    for c in candidates:
        if os.path.exists(c):
            return c
    sys.exit("master not found (looked in: %s)" % ", ".join(candidates))


def placeholder_rows(layout):
    rows = []
    for ph in layout.placeholders:
        f = ph.placeholder_format
        rows.append({
            "idx": f.idx,
            "type": str(f.type).split()[0] if f.type is not None else "None",
            "name": ph.name,
        })
    return rows


def main():
    as_json = "--json" in sys.argv[1:]
    master = resolve_master()
    prs = Presentation(master)
    layouts = prs.slide_masters[0].slide_layouts

    data = []
    for i, layout in enumerate(layouts):
        data.append({
            "index": i,
            "name": layout.name,
            "placeholders": placeholder_rows(layout),
        })

    if as_json:
        print(json.dumps({"master": os.path.relpath(master, REPO_ROOT),
                          "layoutCount": len(data), "layouts": data}, indent=2))
        return

    print("Master: %s" % os.path.relpath(master, REPO_ROOT))
    print("Layouts: %d\n" % len(data))
    for d in data:
        phs = ", ".join("%d:%s(%s)" % (p["idx"], p["name"], p["type"]) for p in d["placeholders"]) or "(none)"
        print("[%2d] %-22s  %s" % (d["index"], d["name"], phs))


if __name__ == "__main__":
    main()
